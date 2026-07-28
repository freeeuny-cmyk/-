import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

target_files = [
    r"c:\Users\user\Desktop\프로젝트(숏폼)\숏폼 스튜디어(기획교육과_이동은).pptx",
    r"C:\Users\user\Desktop\숏폼스튜디오(기획교육과_이동은).pptx"
]

url = "https://gban-shorts.onrender.com/"

for filepath in target_files:
    if not os.path.exists(filepath):
        print(f"Skipping (not found): {filepath}")
        continue
    
    print(f"Processing: {filepath}")
    prs = Presentation(filepath)
    last_slide = prs.slides[-1]
    
    # Check if there is an existing text box with onrender text
    target_textbox = None
    for shape in last_slide.shapes:
        if shape.has_text_frame and "onrender" in shape.text_frame.text.lower():
            target_textbox = shape
            break
            
    if target_textbox:
        tf = target_textbox.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = url
        run.font.name = "Malgun Gothic"
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.underline = True
        run.font.color.rgb = RGBColor(16, 185, 129)
        run.hyperlink.address = url
    else:
        # Create text box for URL
        tb = last_slide.shapes.add_textbox(Inches(4.64), Inches(8.88), Inches(8.5), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = url
        run.font.name = "Malgun Gothic"
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.underline = True
        run.font.color.rgb = RGBColor(16, 185, 129)
        run.hyperlink.address = url

    # Remove existing button if already created
    for shape in list(last_slide.shapes):
        if shape.has_text_frame and "바로가기" in shape.text_frame.text:
            sp_element = shape.element
            sp_element.getparent().remove(sp_element)

    # Add dedicated Clickable Button Shape
    btn = last_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.64), Inches(7.85), Inches(8.5), Inches(0.85))
    btn.fill.solid()
    btn.fill.fore_color.rgb = RGBColor(16, 185, 129) # Emerald Green
    btn.line.color.rgb = RGBColor(255, 255, 255)
    btn.line.width = Pt(1.5)
    btn.click_action.hyperlink.address = url
    
    btf = btn.text_frame
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    bp.text = "🚀 GBAN AI 숏폼 퀵스튜디오 바로가기 (클릭) 🔗"
    bp.font.name = "Malgun Gothic"
    bp.font.size = Pt(18)
    bp.font.bold = True
    bp.font.color.rgb = RGBColor(255, 255, 255)
    
    prs.save(filepath)
    print(f"Successfully updated and saved: {filepath}")

print("All PPTX files updated successfully!")
