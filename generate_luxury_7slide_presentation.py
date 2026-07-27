import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import win32com.client
import fitz

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Premium Luxury Color Palette (AI_Agricultural_Disaster_Shorts_Studio Palette)
    C_BG = RGBColor(246, 248, 245)        # Soft Luxury Sand #f6f8f5
    C_CARD = RGBColor(255, 255, 255)      # Pure White Card #ffffff
    C_DARK = RGBColor(19, 34, 27)         # Deep Forest Dark #13221b
    C_CARD_DARK = RGBColor(27, 46, 37)    # Subdued Dark Card #1b2e25
    C_PRIMARY = RGBColor(30, 64, 48)      # Deep Sage Green #1e4030
    C_ACCENT = RGBColor(16, 185, 129)     # Emerald Green #10b981
    C_ACCENT_LIGHT = RGBColor(167, 243, 208) # Light Mint #a7f3d0
    C_TEXT_DARK = RGBColor(20, 30, 24)    # Charcoal Text #141e18
    C_MUTED = RGBColor(85, 98, 88)       # Olive Gray #556258
    C_BORDER = RGBColor(226, 232, 224)    # Subtle Border #e2e8e0
    C_RED_ACCENT = RGBColor(225, 29, 72)  # Rose Red #e11d48
    C_RED_BG = RGBColor(255, 241, 242)    # Light Red Tint

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_BG
        bg.line.fill.background()

    def add_card(slide, left, top, width, height, bg_color=C_CARD, border_color=C_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()
        return card

    def add_header(slide, slide_num_str, cat_text, title_text):
        add_bg(slide)
        
        # Header Graphic Line / Bar Accent
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(0.12), Inches(0.9))
        bar.fill.solid()
        bar.fill.fore_color.rgb = C_ACCENT
        bar.line.fill.background()

        # Category Badge Text
        cat_box = slide.shapes.add_textbox(Inches(1.05), Inches(0.35), Inches(11.4), Inches(0.35))
        p = cat_box.text_frame.paragraphs[0]
        p.text = f"{slide_num_str}  |  {cat_text}"
        p.font.name = "Outfit"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_ACCENT

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(1.05), Inches(0.68), Inches(11.4), Inches(0.7))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_DARK

    # =========================================================================
    # SLIDE 1: Cover (표지 - AI_Agricultural_Disaster_Shorts_Studio 스타일)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_DARK
    bg1.line.fill.background()

    # Decorative Accents
    dec = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    dec.fill.solid()
    dec.fill.fore_color.rgb = C_CARD_DARK
    dec.line.color.rgb = RGBColor(38, 64, 52)
    dec.line.width = Pt(1.5)

    # Category Pill Badge
    pill = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(1.3), Inches(7.5), Inches(0.45))
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(30, 64, 48)
    pill.line.color.rgb = C_ACCENT
    pill.line.width = Pt(1)

    tb = s1.shapes.add_textbox(Inches(1.4), Inches(1.32), Inches(7.3), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = "경상북도농업기술원 AI 데이터 전문인재 양성 프로젝트 발표"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT_LIGHT

    # Main Title
    tb = s1.shapes.add_textbox(Inches(1.3), Inches(2.0), Inches(10.7), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "농업재해 및 병해충 긴급 전파를 위한\nAI 숏폼(Shorts) 자동 생성기"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    tb = s1.shapes.add_textbox(Inches(1.3), Inches(4.0), Inches(10.7), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "부제: 기존 텍스트 문자의 한계를 극복하는 고가독성 재해 대응 속보 제작 플랫폼"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(17)
    p.font.color.rgb = RGBColor(167, 243, 208)

    # Footer Card Info
    fc = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(5.1), Inches(10.733), Inches(1.1))
    fc.fill.solid()
    fc.fill.fore_color.rgb = RGBColor(19, 34, 27)
    fc.line.color.rgb = RGBColor(40, 75, 60)
    fc.line.width = Pt(1)

    tb = s1.shapes.add_textbox(Inches(1.5), Inches(5.25), Inches(10.3), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• 발표팀: 팀 3 (이동은, 김영아, 이미향)\n• 라이브 서비스 주소: https://gban-shorts.onrender.com"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # =========================================================================
    # SLIDE 2: Background & Problem (배경 및 기획 의도)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "SLIDE 02", "BACKGROUND & PROBLEM", "기존 농업재해 전파 방식의 한계와 기획 배경")

    # Left Card - Problem (Red Accent)
    add_card(s2, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2), C_CARD, RGBColor(254, 202, 202))
    
    tag1 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(1.9), Inches(5.1), Inches(0.55))
    tag1.fill.solid()
    tag1.fill.fore_color.rgb = C_RED_BG
    tag1.line.color.rgb = C_RED_ACCENT

    tb = s2.shapes.add_textbox(Inches(1.2), Inches(1.95), Inches(4.9), Inches(0.45))
    p = tb.text_frame.paragraphs[0]
    p.text = "🚨 현장 긴급 재난문자(LMS)의 3대 한계"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = C_RED_ACCENT

    tb = s2.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(5.1), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    probs = [
        "1. 가독성 저하: 장문의 텍스트 위주 공지로 농가의 정보 이탈률 극심",
        "2. 고령 농업인 취약: 시각적/청각적 직관성 부족으로 긴급 행동요령 전달 실패",
        "3. 제작 지연: 담당 공무원의 반복적인 공지글 작성 및 동영상 편집 기술 부재"
    ]
    for idx, pr in enumerate(probs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = pr
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(14)
        p.font.color.rgb = C_TEXT_DARK
        p.space_after = Pt(14)

    # Right Card - Solution (Emerald Accent)
    add_card(s2, Inches(6.833), Inches(1.6), Inches(5.7), Inches(5.2), C_CARD, C_ACCENT)
    
    tag2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.133), Inches(1.9), Inches(5.1), Inches(0.55))
    tag2.fill.solid()
    tag2.fill.fore_color.rgb = RGBColor(236, 253, 245)
    tag2.line.color.rgb = C_ACCENT

    tb = s2.shapes.add_textbox(Inches(7.233), Inches(1.95), Inches(4.9), Inches(0.45))
    p = tb.text_frame.paragraphs[0]
    p.text = "💡 AI 숏폼 스튜디오의 혁신적 해결책"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT

    tb = s2.shapes.add_textbox(Inches(7.133), Inches(2.6), Inches(5.1), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    sols = [
        "✔ 1~2분 초스피드 제작: 사진 첨부 + 문구 입력만으로 즉시 동영상 자동 생성",
        "✔ 융합형 숏폼 포맷: 자막 + 고품질 AI 음성(TTS) + BGM이 조화된 9:16 비디오",
        "✔ 대민 신속 전파: 완성된 MP4 동영상을 카카오톡, 재난문자, SNS로 즉시 보급"
    ]
    for idx, so in enumerate(sols):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = so
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(14)
        p.font.color.rgb = C_TEXT_DARK
        p.space_after = Pt(14)

    # =========================================================================
    # SLIDE 3: Key Features (핵심 주요 기능 4가지 Grid Layout)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "SLIDE 03", "KEY FEATURES", "웹앱의 4대 핵심 서비스 기능")

    card_w = Inches(5.66)
    card_h = Inches(2.45)
    positions = [
        (Inches(0.8), Inches(1.6)),
        (Inches(6.86), Inches(1.6)),
        (Inches(0.8), Inches(4.35)),
        (Inches(6.86), Inches(4.35))
    ]
    f_data = [
        ("01", "드래그 앤 드롭 사진 매핑", "현장 피해 사진 및 대응 순서 이미지를 손쉽게 업로드하고 드래그/화살표 버튼으로 자막과 1:1 매칭"),
        ("02", "고품질 한국어 AI 음성 (TTS)", "OpenAI 시머, 노바, 오닉스 등 최고급 성우 목소리 지원 및 음성 재생 속도(0.5x~2.0x) 세밀 조절"),
        ("03", "농정 전용 BGM 오디오 믹싱", "Fields of Opportunity 등 고품질 배경음악 5종 탑재, 실시간 미리듣기 및 음성/BGM 믹싱 볼륨 제어"),
        ("04", "9:16 모바일 MP4 동영상 생성", "경북 전용 워터마크 자동 탑재, 캔버스 2D 라이브 렌더링 후 표준 MP4 다운로드 & QR코드 모바일 접속")
    ]
    for idx, (num, title, desc) in enumerate(f_data):
        lx, ty = positions[idx]
        add_card(s3, lx, ty, card_w, card_h, C_CARD, C_BORDER)
        
        # Step Circle/Badge
        nb = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx + Inches(0.3), ty + Inches(0.3), Inches(0.7), Inches(0.4))
        nb.fill.solid()
        nb.fill.fore_color.rgb = C_PRIMARY
        nb.line.fill.background()

        tb = s3.shapes.add_textbox(lx + Inches(0.3), ty + Inches(0.28), Inches(0.7), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = num
        p.font.name = "Outfit"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Title
        tb = s3.shapes.add_textbox(lx + Inches(1.15), ty + Inches(0.26), card_w - Inches(1.3), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_DARK

        # Desc
        tb = s3.shapes.add_textbox(lx + Inches(0.3), ty + Inches(0.85), card_w - Inches(0.6), card_h - Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(14)
        p.font.color.rgb = C_MUTED

    # =========================================================================
    # SLIDE 4: User Workflow (사용자 이용 프로세스 4 Step Cards)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "SLIDE 04", "SERVICE WORKFLOW", "현장 담당자 4단계 이용 프로세스")

    step_w = Inches(2.68)
    gap = Inches(0.33)
    steps = [
        ("STEP 01", "사진 첨부", "농업재해 피해 사진 또는 작물 대응 가이드 이미지를 순서대로 업로드"),
        ("STEP 02", "대본 작성", "줄바꿈(엔터) 단위로 자막 문구 작성, 슬라이드별 재생시간 및 위치 설정"),
        ("STEP 03", "음성 & BGM", "원하는 AI 보이스 선택, 5종 농정 BGM 및 오디오 사운드 미리듣기 확인"),
        ("STEP 04", "동영상 다운로드", "브라우저 캔버스 자동 생성 후 MP4 다운로드 및 카카오톡/문자 즉시 전파")
    ]
    for idx, (s_num, s_title, s_desc) in enumerate(steps):
        lx = Inches(0.8) + idx * (step_w + gap)
        add_card(s4, lx, Inches(1.6), step_w, Inches(5.2), C_CARD, C_BORDER)
        
        # Step Header Pill
        spill = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx + Inches(0.2), Inches(1.9), step_w - Inches(0.4), Inches(0.5))
        spill.fill.solid()
        spill.fill.fore_color.rgb = C_PRIMARY
        spill.line.fill.background()

        tb = s4.shapes.add_textbox(lx + Inches(0.2), Inches(1.95), step_w - Inches(0.4), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = s_num
        p.font.name = "Outfit"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = C_ACCENT_LIGHT

        # Title
        tb = s4.shapes.add_textbox(lx + Inches(0.2), Inches(2.6), step_w - Inches(0.4), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s_title
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(19)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_DARK

        # Desc
        tb = s4.shapes.add_textbox(lx + Inches(0.2), Inches(3.5), step_w - Inches(0.4), Inches(3.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s_desc
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(14)
        p.font.color.rgb = C_MUTED

    # =========================================================================
    # SLIDE 5: Tech Architecture (기술 스택 및 아키텍처)
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "SLIDE 05", "TECH ARCHITECTURE", "웹앱 기술 아키텍처 및 웹 표준 스택")

    col_w2 = Inches(5.7)
    add_card(s5, Inches(0.8), Inches(1.6), col_w2, Inches(5.2), C_CARD, C_BORDER)
    
    tb = s5.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.1), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "💻 Frontend & Web Standards"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    tb = s5.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(5.1), Inches(4.1))
    tf = tb.text_frame
    tf.word_wrap = True
    
    front_techs = [
        "• Vanilla JavaScript: 외부 프레임워크 없는 초고속 반응속도",
        "• Canvas 2D API: 9:16 모바일 비디오 실시간 그래픽 렌더링",
        "• Web Audio API: 음성+BGM 실시간 오디오 스트리밍 & 믹싱",
        "• MediaRecorder API: 웹 브라우저 단독 MP4 비디오 인코딩",
        "• GOBE Luxury System: 다크모드 반전 방지 라이트모드 UI"
    ]
    for idx, ft in enumerate(front_techs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ft
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(14)
        p.font.color.rgb = C_TEXT_DARK
        p.space_after = Pt(10)

    add_card(s5, Inches(6.833), Inches(1.6), col_w2, Inches(5.2), C_CARD, C_BORDER)
    
    tb = s5.shapes.add_textbox(Inches(7.133), Inches(1.9), Inches(5.1), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "⚙️ Backend & Cloud Infrastructure"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    tb = s5.shapes.add_textbox(Inches(7.133), Inches(2.5), Inches(5.1), Inches(4.1))
    tf = tb.text_frame
    tf.word_wrap = True
    
    back_techs = [
        "• Python HTTP API Server: 경량 RESTful backend 서버 구축",
        "• OpenAI Speech REST API: 시머, 오닉스 등 고품질 TTS 연동",
        "• API Key Security: Base64 인코딩 및 환경변수 자동 매핑",
        "• 24/7 Cloud Service: Render Cloud 배포로 모바일 언제나 접속",
        "• QR Code Quick Connect: QR스캔만으로 스마트폰 즉시 실행"
    ]
    for idx, bt in enumerate(back_techs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bt
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(14)
        p.font.color.rgb = C_TEXT_DARK
        p.space_after = Pt(10)

    # =========================================================================
    # SLIDE 6: Expected Impact (기대효과 및 성과)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "SLIDE 06", "EXPECTED IMPACT", "기존 방식 대비 도입 성과 및 기대효과")

    metric_w = Inches(3.64)
    gap = Inches(0.4)
    metrics = [
        ("95% 절감", "영상 제작 시간 단축", "기존 1~2시간 소요되던 편집 작업을 단 1~2분으로 단축하여 재해 골든타임 사수"),
        ("300% 향상", "농가 정보 전달력", "텍스트 대비 시각+음성+자막 융합 콘텐츠로 고령 농업인의 재해 대응 이해도 극대화"),
        ("24/7 보급", "모바일 대민 전파", "QR코드 및 클라우드 URL을 통해 시군 센터 어디서나 365일 즉시 영상 보급")
    ]
    for idx, (m_val, m_title, m_desc) in enumerate(metrics):
        lx = Inches(0.8) + idx * (metric_w + gap)
        add_card(s6, lx, Inches(1.6), metric_w, Inches(5.2), C_CARD, C_BORDER)
        
        # Stat Pill Box
        sp = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx + Inches(0.3), Inches(1.9), metric_w - Inches(0.6), Inches(0.9))
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor(236, 253, 245)
        sp.line.color.rgb = C_ACCENT

        tb = s6.shapes.add_textbox(lx + Inches(0.3), Inches(1.95), metric_w - Inches(0.6), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = m_val
        p.font.name = "Outfit"
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = C_ACCENT

        tb = s6.shapes.add_textbox(lx + Inches(0.3), Inches(3.0), metric_w - Inches(0.6), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = m_title
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(19)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_DARK

        tb = s6.shapes.add_textbox(lx + Inches(0.3), Inches(3.9), metric_w - Inches(0.6), Inches(2.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = m_desc
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(14)
        p.font.color.rgb = C_MUTED

    # =========================================================================
    # SLIDE 7: Conclusion & Roadmap (결론 및 향후 로드맵)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "SLIDE 07", "CONCLUSION & ROADMAP", "프로젝트 최종 결론 및 향후 발전 로드맵")

    hero7 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    hero7.fill.solid()
    hero7.fill.fore_color.rgb = C_DARK
    hero7.line.color.rgb = RGBColor(38, 64, 52)
    hero7.line.width = Pt(1.5)

    tb = s7.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(10.9), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🎯 프로젝트 최종 결론"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT

    summary_items = [
        "1. 농업 재해 및 병해충 상황 발생 시 현장에서 즉시 활용 가능한 실전형 AI 숏폼 생성 플랫폼을 완성하였습니다.",
        "2. 별도 편집 기술 없이 몇 번의 클릭으로 제작 가능한 사용자 친화적 웹 UI/UX 시스템을 구현하였습니다.",
        "3. 향후 발전 로드맵:",
        "    • 자동 기상 특보 데이터 연동 숏폼 자동 발송 체계 구축",
        "    • 외국인 농업 근로자를 위한 다국어(베트남어, 캄보디아어 등) 음성 자막 지원 확대",
        "    • 경상북도농업기술원과 함께하는 스마트 농정 및 도민 안전 실현"
    ]
    for item in summary_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(8)

    # Save Presentations
    out_pptx_desktop = os.path.abspath(os.path.join(os.path.expanduser('~'), 'Desktop', 'GBAN_AI_숏폼_스튜디오_발표자료_7장.pptx'))
    out_pptx_local = os.path.abspath('GBAN_AI_숏폼_스튜디오_발표자료_7장.pptx')
    
    prs.save(out_pptx_local)
    prs.save(out_pptx_desktop)
    print(f"Luxury 7-Slide PPTX created successfully at:\n  {out_pptx_desktop}\n  {out_pptx_local}")
    return out_pptx_local

def convert_to_svg(pptx_path):
    out_dir_desktop = os.path.abspath(os.path.join(os.path.expanduser('~'), 'Desktop', 'GBAN_AI_숏폼_스튜디오_발표자료_SVG_7장'))
    out_dir_local = os.path.abspath('public/svg_presentation_7slides')
    pdf_path = os.path.abspath('temp_presentation_7slides_luxury.pdf')

    os.makedirs(out_dir_desktop, exist_ok=True)
    os.makedirs(out_dir_local, exist_ok=True)

    print("\nConverting Luxury 7-Slide PPTX to PDF using PowerPoint COM...")
    app = win32com.client.Dispatch('PowerPoint.Application')
    pres = app.Presentations.Open(pptx_path, WithWindow=False)
    pres.SaveAs(pdf_path, 32)
    pres.Close()

    print("Converting PDF to 7 High-End Vector SVG Slides...")
    doc = fitz.open(pdf_path)
    for i, page in enumerate(doc):
        slide_num = i + 1
        svg_text = page.get_svg_image()
        
        filename = f"slide_presentation_{slide_num:02d}.svg"
        path_local = os.path.join(out_dir_local, filename)
        path_desktop = os.path.join(out_dir_desktop, filename)
        
        with open(path_local, 'w', encoding='utf-8') as f:
            f.write(svg_text)
            
        with open(path_desktop, 'w', encoding='utf-8') as f:
            f.write(svg_text)
            
        print(f"   [Slide {slide_num:02d}/7] -> Saved {filename}")

    doc.close()
    if os.path.exists(pdf_path):
        try: os.remove(pdf_path)
        except Exception: pass

    print("\n==================================================")
    print("  Luxury 7-Slide Masterpiece Presentation Created Successfully!")
    print(f"  Desktop PPTX: {os.path.join(os.path.expanduser('~'), 'Desktop', 'GBAN_AI_숏폼_스튜디오_발표자료_7장.pptx')}")
    print(f"  Desktop SVG Folder: {out_dir_desktop}")
    print("==================================================")

def main():
    local_pptx = build_presentation()
    convert_to_svg(local_pptx)

if __name__ == '__main__':
    main()
