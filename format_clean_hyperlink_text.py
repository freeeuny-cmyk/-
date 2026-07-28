import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

target_files = [
    r"c:\Users\user\Desktop\프로젝트(숏폼)\숏폼 스튜디어(기획교육과_이동은).pptx",
    r"C:\Users\user\Desktop\숏폼스튜디오(기획교육과_이동은).pptx"
]

url = "https://gban-shorts.onrender.com/"

for filepath in target_files:
    if not os.path.exists(filepath):
        print(f"File not found: {filepath}")
        continue
        
    print(f"Processing: {filepath}")
    prs = Presentation(filepath)
    last_slide = prs.slides[-1]
    
    # 1. Remove any AUTO_SHAPE button added previously
    for shape in list(last_slide.shapes):
        if shape.has_text_frame and "바로가기" in shape.text_frame.text:
            sp_elem = shape.element
            sp_elem.getparent().remove(sp_elem)
            
    # 2. Find or create the target URL text box
    url_textbox = None
    for shape in last_slide.shapes:
        if shape.has_text_frame and "onrender" in shape.text_frame.text.lower():
            url_textbox = shape
            break
            
    if not url_textbox:
        url_textbox = last_slide.shapes.add_textbox(Inches(4.0), Inches(8.5), Inches(9.78), Inches(1.0))
    else:
        # Reposition to clean centered bottom position
        url_textbox.left = Inches(4.0)
        url_textbox.top = Inches(8.5)
        url_textbox.width = Inches(9.78)
        url_textbox.height = Inches(1.0)

    # 3. Apply clean underlined hyperlink text formatting
    tf = url_textbox.text_frame
    tf.word_wrap = True
    tf.clear()
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    run = p.add_run()
    run.text = url
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.underline = True
    # Classic Hyperlink Blue #0563C1
    run.font.color.rgb = RGBColor(5, 99, 193)
    # Hyperlink address
    run.hyperlink.address = url
    
    prs.save(filepath)
    print(f"Successfully formatted underline hyperlink text in: {filepath}")

print("All target PPTX files updated with underline hyperlink text!")
