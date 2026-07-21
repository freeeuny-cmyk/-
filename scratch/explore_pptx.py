from pptx import Presentation
import os

pptx_name = "★ 팀프로젝트 기획 방향 양식_공유.pptx"
prs = Presentation(pptx_name)

print(f"Number of slides: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    print(f"\n=== Slide {i+1} ===")
    for j, shape in enumerate(slide.shapes):
        print(f"Shape {j}: Name={shape.name}, Type={shape.shape_type}")
        if shape.has_text_frame and shape.text:
            print(f"  Text: {shape.text}")
        if shape.has_table:
            print(f"  Table structure: {len(shape.table.rows)} rows x {len(shape.table.columns)} columns")
            for r_idx, row in enumerate(shape.table.rows):
                row_texts = []
                for c_idx, cell in enumerate(row.cells):
                    row_texts.append(cell.text.replace('\n', ' '))
                print(f"    Row {r_idx}: {row_texts}")
