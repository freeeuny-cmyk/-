import os
import shutil
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Configure stdout to use UTF-8 encoding
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

pptx_name = "★ 팀프로젝트 기획 방향 양식_공유.pptx"
backup_name = "★ 팀프로젝트 기획 방향 양식_공유_Backup.pptx"

# Backup first if backup doesn't exist
if not os.path.exists(backup_name):
    shutil.copy(pptx_name, backup_name)
    print(f"Backup created at: {backup_name}")

prs = Presentation(pptx_name)
slide = prs.slides[0]

# Find the Table shape (표 8)
table_shape = None
for shape in slide.shapes:
    if shape.has_table:
        table_shape = shape
        break

if not table_shape:
    print("Table not found in slide!")
    sys.exit(1)

table = table_shape.table
print("Found table! Formatting content...")

def set_cell_lines(cell, lines, font_size_pt=10, bold=False):
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = '맑은 고딕'
        p.font.size = Pt(font_size_pt)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(30, 30, 30)  # Dark gray color for neat look

# 1. Row 0: Team representative and members
# Row 0: ['팀 대표 ', '000', '팀원', '000, 000, 000, 000']
# We assume Team 5: 고정우, 이종필, 구은영, 허수현
# We will set Representative to '고정우' and members to '이종필, 구은영, 허수현'
set_cell_lines(table.cell(0, 1), ["고정우"], font_size_pt=11, bold=True)
set_cell_lines(table.cell(0, 3), ["이종필, 구은영, 허수현"], font_size_pt=11, bold=True)

# 2. Row 2: 프로젝트 주제
# '프로젝트 주제' cell(2, 0) -> 내용 cell(2, 1)
project_title = [
    "농업재해 및 병해충 긴급 전파를 위한 기술원 내부용 AI 숏츠(Shorts) 자동 생성기",
    "(부제: 기존 텍스트 문자의 한계를 극복하는 고가독성 재해 대응 속보 제작 도구)"
]
set_cell_lines(table.cell(2, 1), project_title, font_size_pt=11, bold=True)

# 3. Row 3: 프로젝트 목적
project_purpose = [
    "1. 가독성이 떨어지는 기존 텍스트 위주 농업 재해/병해충 안내 문자를 숏츠 영상으로 대체하여 정보 전달력 극대화",
    "2. 고령 농업인 등 시각 및 디지털 취약층도 쉽게 이해할 수 있도록 이미지, 음성(TTS), 자막이 융합된 직관적 콘텐츠 제공",
    "3. 재해 및 병해충 긴급 상황 발생 시 기술원 담당자가 별도 편집 기술 없이 1~2분 내에 대민 전파용 숏츠 영상을 즉시 제작 및 보급"
]
set_cell_lines(table.cell(3, 1), project_purpose, font_size_pt=10)

# 4. Row 4: 사용자 유형
user_types = [
    "• 경상북도농업기술원 및 시군 농업기술센터의 재해 대비, 병해충 예찰, 기술보급 담당 지도사 및 홍보 연구원",
    "• (※ 농업인에게 웹앱을 직접 배포하거나 공유하지 않고, 기술원에서 제작한 결과물 영상만 메신저/SNS/홈페이지로 배포)"
]
set_cell_lines(table.cell(4, 1), user_types, font_size_pt=10)

# 5. Row 5: 반복 행동
repeated_actions = [
    "• 농업재해(태풍, 냉해 등)나 병해충 발생 시, 상황 공유 및 농가 행동 요령안 공지글을 반복적으로 작성함",
    "• 현장 및 작물 피해 사진 등을 수집하고 이를 농민들에게 신속히 전파하기 위해 긴급 문자 메시지(LMS) 발송 대기함",
    "• 홈페이지 공지사항 등록 및 대민 발송용 콘텐츠 포맷팅 작업을 재해 발생 시마다 반복함"
]
set_cell_lines(table.cell(5, 1), repeated_actions, font_size_pt=10)

# 6. Row 6: 도구/환경
tools_env = [
    "• 도구: 재해/병해충 예보 텍스트 매뉴얼, 스마트폰 촬영 사진, 기술원 대민 문자 발송 시스템, 홈페이지 관리 도구",
    "• 환경: 긴급 재해 경보가 발령되는 급박한 상황 속에서 신속하고 즉각적인 정보 제작 및 발송이 요구되는 기술원 사무실 환경"
]
set_cell_lines(table.cell(6, 1), tools_env, font_size_pt=10)

# 7. Row 7: 구체 결핍
detailed_needs = [
    "• 문자 메시지의 가독성 한계: 텍스트 위주의 문자는 내용이 길어 가독성이 낮고, 특히 고령 농민들의 행동 요령 숙지율이 떨어짐",
    "• 오디오/시각 정보 결핍: 단순 텍스트보다 음성으로 설명해 주는 것이 정보 전달 효과가 크나, 매번 음성 녹음과 비디오 편집을 하기에는 시간과 장비가 부족함",
    "• 골든타임 사수 불가: 기존 비디오 편집 도구로는 신속한 상황 전파가 생명인 재해/병해충 발령 속도에 맞춰 영상을 제작할 수 없음"
]
set_cell_lines(table.cell(7, 1), detailed_needs, font_size_pt=10)

# 8. Row 8: 현재 해결방식
current_solution = [
    "• 가독성이 낮은 장문의 LMS 문자 메시지를 농가에 일괄 발송하는 데 그침",
    "• 기술원 홈페이지 공지사항에 정적인 한글 파일이나 이미지(카드뉴스)를 첨부하여 게시함",
    "• 간혹 전문 영상 제작을 시도하더라도 기획부터 편집까지 며칠씩 소요되어 실시간 재해 정보로서의 시의성을 상실함"
]
set_cell_lines(table.cell(8, 1), current_solution, font_size_pt=10)

# 9. Row 9: NEW 해결방안
new_solution = [
    "• 기술원 내부 전용 긴급 숏츠 제작기: 재해/병해충 현장 사진과 대응 요령 텍스트만 넣으면 즉시 숏츠 영상으로 인코딩",
    "• 무편집 타임라인 동기화 엔진: 마침표나 줄바꿈 단위로 자막을 쪼개어 TTS 음성 재생 시간에 맞춰 사진 슬라이드 및 자막 노출 구간을 100% 자동 동기화 렌더링",
    "• 고가독성 모바일 최적화 포맷: 농민들이 스마트폰 메신저나 문자로 받아보았을 때 한눈에 들어오는 세로형 숏츠 영상(MP4) 즉시 다운로드 및 기술원 홈페이지/SNS 채널 게시 연동"
]
set_cell_lines(table.cell(9, 1), new_solution, font_size_pt=10)

# 10. Row 10: 활용 데이터
utilized_data = [
    "• 사용자 입력 데이터: 재해/병해충 피해 사진, 긴급 방제법 및 행동 요령 대본 텍스트",
    "• 시스템 제공 데이터: 경고 및 주의 환기용 긴장감 있는 BGM 에셋, 템플릿 레이아웃 디자인",
    "• AI API 데이터: 자연스러운 경보 및 안내 방송 톤의 OpenAI TTS API 오디오 데이터(onyx, echo 등 신뢰감 있는 남성 톤 선호) 및 Google TTS 데이터"
]
set_cell_lines(table.cell(10, 1), utilized_data, font_size_pt=10)

# 11. Row 11: AI 개입 포인트
ai_intervention = [
    "• AI 기반 긴급 안내 방송 음성 합성: OpenAI TTS의 신뢰감 있는 보이스(예: onyx)를 적용하여 농업 재해 방송의 격식과 가독성을 확보",
    "• 자막 및 시각 연출의 자동 스케줄링: 합성된 재해 경보 TTS 음성의 실시간 재생 길이를 분석하여, 텍스트 자막 오버레이와 재해 상황 이미지의 슬라이드 싱크를 인위적 편집 과정 없이 완전 자동 매칭함"
]
set_cell_lines(table.cell(11, 1), ai_intervention, font_size_pt=10)

output_pptx_name = "★ 팀프로젝트 기획 방향 양식_작성완료_최종.pptx"
prs.save(output_pptx_name)
print(f"Successfully updated pptx file: {output_pptx_name}")
