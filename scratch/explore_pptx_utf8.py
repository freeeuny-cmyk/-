import sys
from pptx import Presentation

# Configure stdout to use UTF-8 encoding
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

pptx_name = "★ 팀프로젝트 기획 방향 양식_공유.pptx"
prs = Presentation(pptx_name)

print(f"Number of slides: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    print(f"\n=== Slide {i+1} ===")
    for j, shape in enumerate(slide.shapes):
        print(f"Shape {j}: Name={shape.name}, Type={shape.shape_type}")
        if shape.has_text_frame and shape.text:
            print(f"  Text: {shape.text}")
        if shape.has_table:
            table = shape.table
            print(f"  Table structure: {len(table.rows)} rows x {len(table.columns)} columns")
            for r_idx, row in enumerate(table.rows):
                row_texts = []
                for c_idx, cell in enumerate(row.cells):
                    row_texts.append(cell.text.replace('\n', ' '))
                print(f"    Row {r_idx}: {row_texts}")
