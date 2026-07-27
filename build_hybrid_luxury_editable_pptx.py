import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import win32com.client
import fitz

def create_hybrid_luxury_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Slide texts matching the positions on each slide
    slide_data = [
        {
            # Slide 1: Cover
            "bg": "slide_img_01.png",
            "boxes": [
                {"text": "경상북도농업기술원 AI 데이터 전문인재 양성교육 팀프로젝트", "left": 1.0, "top": 1.2, "w": 11.3, "h": 0.5, "size": 16, "bold": True, "color": RGBColor(16, 185, 129)},
                {"text": "농업재해 및 병해충 긴급 전파를 위한\nAI 숏츠(Shorts) 자동 생성기", "left": 1.0, "top": 1.8, "w": 11.3, "h": 1.8, "size": 34, "bold": True, "color": RGBColor(255, 255, 255)},
                {"text": "부제: 기존 텍스트 문자의 한계를 극복하는 고가독성 재해 대응 속보 제작 도구", "left": 1.0, "top": 3.8, "w": 11.3, "h": 0.8, "size": 18, "bold": False, "color": RGBColor(200, 225, 210)},
                {"text": "팀 대표: 이동은   |   팀원: 김영아, 이미향", "left": 1.0, "top": 5.2, "w": 11.3, "h": 0.8, "size": 18, "bold": True, "color": RGBColor(255, 255, 255)}
            ]
        },
        {
            # Slide 2: Team Status
            "bg": "slide_img_02.png",
            "boxes": [
                {"text": "TEAM STATUS - 팀 구성 현황", "left": 0.8, "top": 0.6, "w": 11.7, "h": 0.5, "size": 16, "bold": True, "color": RGBColor(16, 185, 129)},
                {"text": "경상북도농업기술원 AI 데이터 전문인재 양성교육", "left": 0.8, "top": 1.0, "w": 11.7, "h": 0.6, "size": 24, "bold": True, "color": RGBColor(26, 46, 38)},
                {"text": "팀 3: 이동은 (팀장), 김영아, 이미향", "left": 0.8, "top": 2.2, "w": 11.7, "h": 1.0, "size": 20, "bold": True, "color": RGBColor(30, 38, 31)},
                {"text": "농업재해 및 병해충 긴급 전파를 위한 AI 숏츠 자동 생성기 개발", "left": 0.8, "top": 3.4, "w": 11.7, "h": 1.0, "size": 18, "bold": False, "color": RGBColor(90, 102, 91)}
            ]
        },
        {
            # Slide 3: Timeline
            "bg": "slide_img_03.png",
            "boxes": [
                {"text": "TIMELINE - 팀 프로젝트 일정 및 방향", "left": 0.8, "top": 0.6, "w": 11.7, "h": 0.5, "size": 16, "bold": True, "color": RGBColor(16, 185, 129)},
                {"text": "8회차, 9회차, 10회차 중심의 단계별 과제 및 고도화", "left": 0.8, "top": 1.0, "w": 11.7, "h": 0.6, "size": 24, "bold": True, "color": RGBColor(26, 46, 38)},
                {"text": "• 8회차 (7/8): 프로젝트 기획 및 주제 구체화\n• 9회차 (7/22): 실전 개발 및 데이터/코드 조율 (80% 달성)\n• 10회차 (7/29): 최종 발표 및 시상", "left": 0.8, "top": 2.2, "w": 11.7, "h": 4.0, "size": 18, "bold": False, "color": RGBColor(30, 38, 31)}
            ]
        },
        {
            # Slide 4: Timeline Detail
            "bg": "slide_img_04.png",
            "boxes": [
                {"text": "TIMELINE DETAIL - 상세 일정", "left": 0.8, "top": 0.6, "w": 11.7, "h": 0.5, "size": 16, "bold": True, "color": RGBColor(16, 185, 129)},
                {"text": "프로젝트 기획, 개발, 최종 고도화 로드맵", "left": 0.8, "top": 1.0, "w": 11.7, "h": 0.6, "size": 24, "bold": True, "color": RGBColor(26, 46, 38)}
            ]
        },
        {
            # Slide 5: Problem Definition Template
            "bg": "slide_img_05.png",
            "boxes": [
                {"text": "문제정의 템플릿 - STEP 1. 사용자 행동 기반 결핍 정의", "left": 0.8, "top": 0.6, "w": 11.7, "h": 0.5, "size": 16, "bold": True, "color": RGBColor(16, 185, 129)},
                {"text": "기존 텍스트 위주 재해 안내의 한계 극복", "left": 0.8, "top": 1.0, "w": 11.7, "h": 0.6, "size": 24, "bold": True, "color": RGBColor(26, 46, 38)},
                {"text": "프로젝트 목적:\n1. 텍스트 위주 긴급문자를 숏츠 영상으로 대체하여 전달력 극대화\n2. 고령 농업인 등 취약계층용 이미지+자막+음성(TTS) 제공\n3. 담당자가 1~2분 만에 대민 전파 숏츠 즉시 제작", "left": 0.8, "top": 2.0, "w": 11.7, "h": 4.5, "size": 17, "bold": False, "color": RGBColor(30, 38, 31)}
            ]
        },
        {
            # Slide 6 ~ 11
            "bg": "slide_img_06.png",
            "boxes": [
                {"text": "STEP 2. 결핍 해결 솔루션 구체화", "left": 0.8, "top": 0.6, "w": 11.7, "h": 0.5, "size": 16, "bold": True, "color": RGBColor(16, 185, 129)},
                {"text": "AI 기반 숏폼 동영상 자동 생성 솔루션", "left": 0.8, "top": 1.0, "w": 11.7, "h": 0.6, "size": 24, "bold": True, "color": RGBColor(26, 46, 38)}
            ]
        },
        {
            "bg": "slide_img_07.png",
            "boxes": [
                {"text": "STEP 3. 핵심 기능 구현 및 연동", "left": 0.8, "top": 0.6, "w": 11.7, "h": 0.5, "size": 16, "bold": True, "color": RGBColor(16, 185, 129)},
                {"text": "TTS 음성 합성 & BGM 오디오 믹싱 기술", "left": 0.8, "top": 1.0, "w": 11.7, "h": 0.6, "size": 24, "bold": True, "color": RGBColor(26, 46, 38)}
            ]
        },
        {
            "bg": "slide_img_08.png",
            "boxes": [
                {"text": "STEP 4. 캔버스 렌더링 및 MP4 생성", "left": 0.8, "top": 0.6, "w": 11.7, "h": 0.5, "size": 16, "bold": True, "color": RGBColor(16, 185, 129)},
                {"text": "웹 브라우저 기반 실시간 숏폼 동영상 다운로드", "left": 0.8, "top": 1.0, "w": 11.7, "h": 0.6, "size": 24, "bold": True, "color": RGBColor(26, 46, 38)}
            ]
        },
        {
            "bg": "slide_img_09.png",
            "boxes": [
                {"text": "STEP 5. 기대효과 및 정보 전달력 비교", "left": 0.8, "top": 0.6, "w": 11.7, "h": 0.5, "size": 16, "bold": True, "color": RGBColor(16, 185, 129)},
                {"text": "농가 재해 대응 속도 및 피해 예방률 극대화", "left": 0.8, "top": 1.0, "w": 11.7, "h": 0.6, "size": 24, "bold": True, "color": RGBColor(26, 46, 38)}
            ]
        },
        {
            "bg": "slide_img_10.png",
            "boxes": [
                {"text": "STEP 6. 현장 파급력 및 발전 방향", "left": 0.8, "top": 0.6, "w": 11.7, "h": 0.5, "size": 16, "bold": True, "color": RGBColor(16, 185, 129)},
                {"text": "스마트 농정 강화를 위한 AI 기술 실증", "left": 0.8, "top": 1.0, "w": 11.7, "h": 0.6, "size": 24, "bold": True, "color": RGBColor(26, 46, 38)}
            ]
        },
        {
            "bg": "slide_img_11.png",
            "boxes": [
                {"text": "THANK YOU - 경상북도농업기술원", "left": 0.8, "top": 0.6, "w": 11.7, "h": 0.5, "size": 16, "bold": True, "color": RGBColor(16, 185, 129)},
                {"text": "농업혁신을 선도하는 데이터 시너지와 AI 기술실증", "left": 0.8, "top": 1.0, "w": 11.7, "h": 0.6, "size": 24, "bold": True, "color": RGBColor(26, 46, 38)}
            ]
        }
    ]

    for idx, sdata in enumerate(slide_data):
        slide = prs.slides.add_slide(blank_layout)
        bg_file = sdata["bg"]
        if os.path.exists(bg_file):
            slide.shapes.add_picture(bg_file, 0, 0, prs.slide_width, prs.slide_height)
            
        for box in sdata["boxes"]:
            tb = slide.shapes.add_textbox(Inches(box["left"]), Inches(box["top"]), Inches(box["w"]), Inches(box["h"]))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = box["text"]
            p.font.name = "Malgun Gothic"
            p.font.size = Pt(box["size"])
            p.font.bold = box["bold"]
            p.font.color.rgb = box["color"]

    output_pptx = os.path.abspath('AI_Agricultural_Disaster_Shorts_Studio_프리미엄_수정가능.pptx')
    prs.save(output_pptx)
    print(f"Hybrid Luxury PPTX created: {output_pptx}")
    return output_pptx

def convert_pptx_to_hybrid_svg(pptx_path):
    out_dir_local = os.path.abspath('svg_luxury_slides')
    out_dir_desktop = os.path.abspath(os.path.join(os.path.expanduser('~'), 'Desktop', 'AI_Shorts_Studio_프리미엄_SVG'))
    pdf_path = os.path.abspath('temp_luxury.pdf')

    os.makedirs(out_dir_local, exist_ok=True)
    os.makedirs(out_dir_desktop, exist_ok=True)

    print("1. Converting Luxury PPTX to PDF using PowerPoint COM...")
    app = win32com.client.Dispatch('PowerPoint.Application')
    pres = app.Presentations.Open(pptx_path, WithWindow=False)
    pres.SaveAs(pdf_path, 32)
    pres.Close()

    print("2. Converting PDF pages to High-End SVG slides with live text overlays...")
    doc = fitz.open(pdf_path)
    for i, page in enumerate(doc):
        slide_num = i + 1
        svg_text = page.get_svg_image()
        
        filename = f"slide_luxury_{slide_num:02d}.svg"
        path_local = os.path.join(out_dir_local, filename)
        path_desktop = os.path.join(out_dir_desktop, filename)
        
        with open(path_local, 'w', encoding='utf-8') as f:
            f.write(svg_text)
            
        with open(path_desktop, 'w', encoding='utf-8') as f:
            f.write(svg_text)
            
        print(f"   [Slide {slide_num:02d}/11] -> Saved {filename}")

    doc.close()
    if os.path.exists(pdf_path):
        try: os.remove(pdf_path)
        except Exception: pass

    print("\n==================================================")
    print("  Hybrid Luxury PPTX & SVG Files Created Successfully!")
    print(f"  Luxury PPTX File: {pptx_path}")
    print(f"  Desktop SVG Folder: {out_dir_desktop}")
    print("==================================================")

def main():
    pptx_path = create_hybrid_luxury_pptx()
    convert_pptx_to_hybrid_svg(pptx_path)

if __name__ == '__main__':
    main()
