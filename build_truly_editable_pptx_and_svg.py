import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import win32com.client
import fitz

def create_editable_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Colors
    C_BG = RGBColor(248, 246, 240)        # Eco Cream Sand #f8f6f0
    C_CARD = RGBColor(255, 255, 255)      # White Card #ffffff
    C_DARK = RGBColor(26, 46, 38)         # Deep Forest #1a2e26
    C_PRIMARY = RGBColor(75, 99, 78)      # Sage Green #4b634e
    C_ACCENT = RGBColor(16, 185, 129)     # Emerald Green #10b981
    C_TEXT = RGBColor(30, 38, 31)         # Dark Text #1e261f
    C_MUTED = RGBColor(90, 102, 91)       # Muted Gray #5a665b
    C_BORDER = RGBColor(230, 225, 213)    # Sand Border #e6e1d5

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_BG
        bg.line.fill.background()

    def add_header(slide, cat_text, title_text):
        add_bg(slide)
        # Category Badge
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        tf = cat_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[ {cat_text} ]"
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_ACCENT

        # Title Text
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.8))
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.name = "Malgun Gothic"
        p2.font.size = Pt(26)
        p2.font.bold = True
        p2.font.color.rgb = C_DARK

    def add_card(slide, left, top, width, height, bg_color=C_CARD, border_color=C_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()
        return card

    # SLIDE 1: Cover
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)
    
    # Hero Card
    hero = add_card(slide1, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9), C_DARK, None)
    
    # Badge
    b1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.3), Inches(10), Inches(0.5))
    p = b1.text_frame.paragraphs[0]
    p.text = "경상북도농업기술원 AI 데이터 전문인재 양성교육 팀프로젝트"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(16)
    p.font.color.rgb = C_ACCENT
    p.font.bold = True

    # Main Title
    t1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(10.8), Inches(1.8))
    tf = t1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "농업재해 및 병해충 긴급 전파를 위한\nAI 숏츠(Shorts) 자동 생성기"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    s1 = slide1.shapes.add_textbox(Inches(1.2), Inches(3.9), Inches(10.8), Inches(0.8))
    tf = s1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "부제: 기존 텍스트 문자의 한계를 극복하는 고가독성 재해 대응 속보 제작 도구"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(200, 225, 210)

    # Team Info Box
    tinfo = slide1.shapes.add_textbox(Inches(1.2), Inches(5.2), Inches(10.8), Inches(1.0))
    tf = tinfo.text_frame
    p = tf.paragraphs[0]
    p.text = "• 팀 대표: 이동은    • 팀원: 김영아, 이미향"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # SLIDE 2: Problem Definition (결핍 정의)
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "STEP 1. 사용자 행동 기반 결핍 정의", "기존 농업재해 문자(LMS)의 한계와 새로운 해결책")
    
    # Left Card - Problem
    add_card(slide2, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0), C_CARD, C_BORDER)
    tb = slide2.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.1), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🚨 기존 텍스트 긴급문자의 한계"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 38, 38)
    
    bullets = [
        "긴 장문의 텍스트 위주 공지로 가독성 저하",
        "고령 농업인 등 시각/디지털 취약 계층의 정보 이탈",
        "텍스트만으로는 현장 사진 및 행동 요령 직관적 전달 불가능",
        "긴급 상황 발생 시 빠른 전파 효과 미흡"
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(15)
        p.font.color.rgb = C_TEXT

    # Right Card - Solution
    add_card(slide2, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), C_CARD, C_ACCENT)
    tb = slide2.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💡 AI 숏츠 자동 생성기의 혁신"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_ACCENT

    solutions = [
        "사진 + 대본 텍스트 입력만으로 1~2분 만에 숏폼 자동 제작",
        "자막 + 고품질 한국어 AI 음성(TTS) + BGM 복합 제공",
        "고령 농가도 한눈에 이해하는 직관적 9:16 세로형 숏츠 포맷",
        "기술원 담당자가 손쉽게 대민 신속 전파 가능"
    ]
    for s in solutions:
        p = tf.add_paragraph()
        p.text = "✔ " + s
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(15)
        p.font.color.rgb = C_TEXT

    # SLIDE 3: Key Objectives (핵심 목표 3가지)
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "PROJECT GOALS", "프로젝트 3대 핵심 추진 목표")
    
    col_w = Inches(3.64)
    gap = Inches(0.4)
    objs = [
        ("01", "정보 전달력 극대화", "가독성이 떨어지는 기존 텍스트 위주 공지를 숏츠 영상으로 대체하여 농가의 재해 대응 이해도 300% 향상"),
        ("02", "취약계층 접근성 강화", "고령 농업인 등 visual/audio 정보가 필요한 농민에게 이미지+자막+AI음성(TTS)을 동시 제공하여 재해 피해 최소화"),
        ("03", "1~2분 초스피드 제작", "기술원 담당자가 별도 영상 편집 기술 없이 몇 번의 클릭만으로 대민 전파용 숏폼을 즉시 제작 및 배포")
    ]
    for idx, (num, title, desc) in enumerate(objs):
        left_pos = Inches(0.8) + idx * (col_w + gap)
        add_card(slide3, left_pos, Inches(1.8), col_w, Inches(5.0), C_CARD, C_BORDER)
        
        # Num Badge
        nb = slide3.shapes.add_textbox(left_pos + Inches(0.3), Inches(2.1), Inches(3.0), Inches(0.6))
        p = nb.text_frame.paragraphs[0]
        p.text = f"GOAL {num}"
        p.font.name = "Outfit"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = C_ACCENT
        
        # Title
        tb = slide3.shapes.add_textbox(left_pos + Inches(0.3), Inches(2.8), Inches(3.0), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = C_DARK
        
        # Desc
        db = slide3.shapes.add_textbox(left_pos + Inches(0.3), Inches(3.7), Inches(3.0), Inches(2.8))
        tf = db.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(14)
        p.font.color.rgb = C_MUTED

    # SLIDE 4: Target Users (타겟 사용자 분석)
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "TARGET USERS", "타겟 사용자 및 업무 시나리오")
    
    add_card(slide4, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), C_CARD, C_BORDER)
    tb = slide4.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(10.9), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "👤 주요 사용자: 경상북도농업기술원 및 시군 농업기술센터 담당 공무원"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_DARK
    
    lines = [
        "",
        "• 주요 업무: 농업재해(태풍, 냉해, 가뭄 등) 및 병해충 예찰/발생 시 농가 행동요령 긴급 공지",
        "• 기존 페인포인트:",
        "    - 재해 발생 시마다 반복적인 문자 공지글 및 포맷팅 작업 수행",
        "    - 전문 영상 편집 프로그램(프리미어 등) 사용 숙련도 부족으로 즉각적 영상 제작 불가능",
        "• 숏폼 스튜디오 도입 후:",
        "    - 웹 브라우저 접속 후 사진 첨부 + 문구 입력으로 1분 만에 숏폼 동영상 완성",
        "    - 생성된 MP4 영상 파일을 카카오톡, 재난안전 문자, SNS, 홈페이지로 즉시 전파"
    ]
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(16)
        p.font.color.rgb = C_TEXT

    # SLIDE 5: Core Features Overview (핵심 기능 요약)
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "SYSTEM FEATURES", "GBAN AI 숏폼 퀵스튜디오 핵심 기능 4가지")
    
    grid_w = Inches(5.66)
    grid_h = Inches(2.3)
    coords = [
        (Inches(0.8), Inches(1.8)),
        (Inches(6.86), Inches(1.8)),
        (Inches(0.8), Inches(4.4)),
        (Inches(6.86), Inches(4.4))
    ]
    feats = [
        ("01. 드래그앤드롭 사진 업로드", "현장 피해 사진 및 대응 행동요령 순서를 자유롭게 드래그하여 배치 및 변경"),
        ("02. 자동 자막 분할 & 스타일", "엔터(줄바꿈) 단위로 자막 자동 매칭, 자막 위치/색상/크기/재생시간 자율 설정"),
        ("03. 고품질 한국어 AI 음성 (TTS)", "OpenAI 시머/노바/오닉스 등 고품질 유료 목소리 및 무제한 무료 음성 선택 지원"),
        ("04. 9:16 모바일 숏폼 MP4 다운로드", "경북 전용 워터마크 자동 삽입, BGM 오디오 믹싱, 표준 MP4 동영상 렌더링")
    ]
    for idx, (head, text) in enumerate(feats):
        lx, ty = coords[idx]
        add_card(slide5, lx, ty, grid_w, grid_h, C_CARD, C_BORDER)
        tb = slide5.shapes.add_textbox(lx + Inches(0.3), ty + Inches(0.2), grid_w - Inches(0.6), grid_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = head
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = C_ACCENT
        
        p = tf.add_paragraph()
        p.text = text
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(14)
        p.font.color.rgb = C_MUTED

    # SLIDE 6 ~ 11
    slide_titles = [
        ("FEATURE 01", "슬라이드별 자동 자막 및 시간 동기화 기술"),
        ("FEATURE 02", "한국어 AI 음성(TTS) & 배경음악 오디오 믹싱"),
        ("FEATURE 03", "경상북도 워터마크 및 모바일 시뮬레이터 캔버스"),
        ("ARCHITECTURE", "시스템 구조 및 웹표준 기술 스택 (Python/WebAudio/Canvas)"),
        ("BENEFITS", "현장 도입 기대효과 및 정보 전달력 비교"),
        ("SUMMARY", "경상북도농업기술원 AI 숏폼 스튜디오 요약 및 향후 발전 방향")
    ]

    contents = [
        ["• 줄바꿈(엔터) 기준으로 자막과 이미지가 1:1 매칭되는 직관적 슬라이드 매핑", "• 자막 위치(상단/중앙/하단), 글자 색상, 폰트 크기 실시간 변경 가능", "• 슬라이드별 최적 딜레이 계산으로 영상 시각적 안정성 보장"],
        ["• OpenAI TTS (Shimmer, Nova, Alloy, Onyx, Echo, Fable) 연동", "• 5가지 전용 농정 BGM (Fields of Opportunity, 푸른 들녘의 혁신 등) 탑재", "• 음성 속도(0.5x~2.0x) 및 BGM 볼륨 조절 기능 제공"],
        ["• 경상북도 공식 Emblem 상단/하단 워터마크 자동 렌더링", "• 모바일 스마트폰 9:16 종횡비 직관적 실시간 라이브 시뮬레이션", "• 제작 전 음성+BGM 개별/합동 사운드 미리듣기 지원"],
        ["• Backend: Python HTTP Server / FastAPI / OpenAI Speech API Integration", "• Frontend: Vanilla JS, Web Audio API, Canvas 2D, MediaRecorder", "• Cloud: 24/7 Render Cloud Deployment & Mobile Web Standards"],
        ["• 기존 긴급 문자(LMS) 대비 농가 대응 속도 및 정보 이해도 극대화", "• 고령 농업인의 시감각적 정보 접근 장벽 완전 해소", "• 재해 대응 골든타임 내 초스피드 대민 전파 시스템 구축"],
        ["• 농업 현장에 최적화된 맞춤형 AI 숏폼 생성 플랫폼 완성", "• 모바일 QR코드 및 24시간 365일 언제나 접속 가능한 주소 제공", "• 농업 재해 예방 및 도민 안전을 위한 스마트 농정 실현"]
    ]

    for idx, (cat, title) in enumerate(slide_titles):
        slide = prs.slides.add_slide(blank_layout)
        add_header(slide, cat, title)
        add_card(slide, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0), C_CARD, C_BORDER)
        
        tb = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.9), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        for c_idx, line in enumerate(contents[idx]):
            p = tf.paragraphs[0] if c_idx == 0 else tf.add_paragraph()
            p.text = line
            p.font.name = "Malgun Gothic"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = C_DARK
            p.space_after = Pt(20)

    # Save Presentation
    output_pptx = os.path.abspath('AI_Agricultural_Disaster_Shorts_Studio_진짜수정가능.pptx')
    prs.save(output_pptx)
    print(f"Editable PPTX created: {output_pptx}")
    return output_pptx

def convert_pptx_to_editable_svg(pptx_path):
    out_dir_local = os.path.abspath('svg_editable_slides')
    out_dir_desktop = os.path.abspath(os.path.join(os.path.expanduser('~'), 'Desktop', 'AI_Shorts_Studio_진짜수정가능_SVG'))
    pdf_path = os.path.abspath('temp_editable.pdf')

    os.makedirs(out_dir_local, exist_ok=True)
    os.makedirs(out_dir_desktop, exist_ok=True)

    print("1. Converting Editable PPTX to PDF using PowerPoint COM...")
    app = win32com.client.Dispatch('PowerPoint.Application')
    pres = app.Presentations.Open(pptx_path, WithWindow=False)
    pres.SaveAs(pdf_path, 32)
    pres.Close()

    print("2. Converting PDF pages to 100% Vector & Editable SVG slides...")
    doc = fitz.open(pdf_path)
    for i, page in enumerate(doc):
        slide_num = i + 1
        svg_text = page.get_svg_image()
        
        filename = f"slide_editable_{slide_num:02d}.svg"
        path_local = os.path.join(out_dir_local, filename)
        path_desktop = os.path.join(out_dir_desktop, filename)
        
        with open(path_local, 'w', encoding='utf-8') as f:
            f.write(svg_text)
            
        with open(path_desktop, 'w', encoding='utf-8') as f:
            f.write(svg_text)
            
        print(f"   [Slide {slide_num:02d}/11] -> Saved {filename}")

    doc.close()
    if os.path.exists(pdf_path):
        try: os.remove(pdf_path)
        except Exception: pass

    print("\n==================================================")
    print("  Truly Editable PPTX & SVG Files Created Successfully!")
    print(f"  Editable PPTX File: {pptx_path}")
    print(f"  Desktop SVG Folder: {out_dir_desktop}")
    print("==================================================")

def main():
    pptx_path = create_editable_presentation()
    convert_pptx_to_editable_svg(pptx_path)

if __name__ == '__main__':
    main()
