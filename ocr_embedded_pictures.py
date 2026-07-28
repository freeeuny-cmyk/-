import os
import sys
import subprocess
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation('숏폼 스튜디어(기획교육과_이동은).pptx')
print(f"Total slides: {len(prs.slides)}")

def run_ps_ocr(img_path):
    abs_path = os.path.abspath(img_path).replace("'", "''")
    ps_cmd = f'''
    $imgPath = '{abs_path}'
    [Windows.Storage.StorageFile, Windows.Storage, ContentType=WindowsRuntime] | Out-Null
    [Windows.Media.Ocr.OcrEngine, Windows.Foundation.UniversalApiContract, ContentType=WindowsRuntime] | Out-Null
    
    $fileTask = [Windows.Storage.StorageFile]::GetFileFromPathAsync($imgPath)
    $asTask = [System.WindowsRuntimeSystemExtensions].GetMethods() | Where-Object {{ $_.Name -eq 'AsTask' -and $_.GetParameters().Length -eq 1 -and $_.GetParameters()[0].ParameterType.Name.StartsWith('IAsyncOperation') }} | Select-Object -First 1
    
    $file = $asTask.MakeGenericMethod([Windows.Storage.StorageFile]).Invoke($null, @($fileTask)).Result
    $streamTask = $file.OpenAsync([Windows.Storage.FileAccessMode]::Read)
    $stream = $asTask.MakeGenericMethod([Windows.Storage.Streams.IRandomAccessStream]).Invoke($null, @($streamTask)).Result
    
    $decoderTask = [Windows.Graphics.Imaging.BitmapDecoder]::CreateAsync($stream)
    $decoder = $asTask.MakeGenericMethod([Windows.Graphics.Imaging.BitmapDecoder]).Invoke($null, @($decoderTask)).Result
    
    $bmpTask = $decoder.GetSoftwareBitmapAsync()
    $bmp = $asTask.MakeGenericMethod([Windows.Graphics.Imaging.SoftwareBitmap]).Invoke($null, @($bmpTask)).Result
    
    $engine = [Windows.Media.Ocr.OcrEngine]::TryCreateFromUserProfileLanguages()
    $ocrTask = $engine.RecognizeAsync($bmp)
    $ocrResult = $asTask.MakeGenericMethod([Windows.Media.Ocr.OcrResult]).Invoke($null, @($ocrTask)).Result
    
    $ocrResult.Text
    '''
    try:
        res = subprocess.run(["powershell", "-NoProfile", "-Command", ps_cmd], capture_output=True, text=True, encoding='cp949', errors='ignore')
        return res.stdout.strip()
    except Exception as e:
        return str(e)

output_lines = []
for i, slide in enumerate(prs.slides):
    slide_header = f"================ SLIDE {i+1:02d} ================"
    print(slide_header)
    output_lines.append(slide_header)
    for s in slide.shapes:
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_bytes = s.image.blob
            if len(img_bytes) > 50000:  # Main slide background/content pictures
                tmp_fn = f"temp_img_{i+1:02d}.png"
                with open(tmp_fn, 'wb') as f:
                    f.write(img_bytes)
                text = run_ps_ocr(tmp_fn)
                print(text)
                output_lines.append(text)
                if os.path.exists(tmp_fn):
                    os.remove(tmp_fn)

with open("ocr_slide_content.txt", "w", encoding="utf-8") as f:
    f.write("\n".join(output_lines))

print("\nSuccessfully saved OCR slide content to ocr_slide_content.txt!")
